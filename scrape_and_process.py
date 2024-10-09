import requests
from bs4 import BeautifulSoup
import pandas as pd
from pptx.util import Inches
import os

def extract_make_model(make_model_tag):
    if make_model_tag:
        car_model_tag = make_model_tag.find('strong').a
        if car_model_tag:
            make_model = car_model_tag.get_text(strip=True)
            return make_model
    return None

def searchanalysis_makemodel_depre_motorist(url):
    search_analysis = []
    page_number = 1
    na_counter = 0

    while page_number <= 5:
        current_url = f"{url}&page={page_number}"
        response = requests.get(current_url)
        soup = BeautifulSoup(response.text, 'html.parser')

        make_model_divs = soup.find_all('div', class_='make-model')
        for make_model_div in make_model_divs:
            make_model_p = make_model_div.find('p', class_='font-weight-bold mb-0')
            make_model = make_model_p.get_text(strip=True) if make_model_p else 'N.A'
            
            price_div = make_model_div.find_next('div', class_='price-registration-owner')
            price_span = price_div.find('span', class_='text-green') if price_div else None
            price = price_span.get_text(strip=True) if price_span else 'N.A'

            depre_div = make_model_div.find_next('div', class_='depre')
            depreciation = depre_div.get_text(strip=True) if depre_div else 'N.A'

            if price == 'N.A' and depreciation == 'N.A':
                na_counter += 1
            else:
                na_counter = 0

            if na_counter >= 4:
                return search_analysis

            search_analysis.append({'Make-Model': make_model, 'Depreciation': depreciation, 'Price': price})
        
        page_number += 1

    return search_analysis

def searchanalysis_makemodel_depre_sgcar(url):
    search_analysis = {}
    encountered_combinations = set()
    page_number = 0

    while page_number <= 180:
        current_url = f"{url}&BRSR={page_number}"
        response = requests.get(current_url)
        soup = BeautifulSoup(response.text, 'html.parser')

        make_model_tags = soup.find_all('div', style='width:186px;padding-left:4px;', class_='font_13')
        price_tags = soup.find_all('div', style='width:67px; font-weight: 500;', class_='font_12')
        depreciation_listings = soup.find_all('td', width='101')

        for make_model_tag, price_tag, depreciation in zip(make_model_tags, price_tags, depreciation_listings):
            make_model = extract_make_model(make_model_tag) if make_model_tag else None
            price = price_tag.get_text(strip=True).replace(',', '').replace('$', '') if price_tag else 'N.A'
            depreciation_value = depreciation.get_text(strip=True).replace(',', '').replace('/yr', '').strip().replace('$', '') if depreciation else 'N.A'

            combination = (make_model, depreciation_value, price)
            if combination not in encountered_combinations:
                search_analysis[combination] = price
                encountered_combinations.add(combination)

        page_number += 20

    return search_analysis

def write_to_csv_motorist(data, filename, working_directory):  # Pass working directory as a parameter
    df = pd.DataFrame(data, columns=['Make-Model', 'Depreciation', 'Price'])
    df = df[~(((df['Depreciation'] == 'N.A') & (df['Price'] == 'N.A')) | (df['Make-Model'] == 'N.A'))]
    df = df.sort_values(by='Depreciation', key=lambda x: x.replace('N.A', '999999'))

    if not df.empty:
        if not filename.endswith('.csv'):
            filename += '.csv'
        filename = os.path.join(working_directory, filename)  # Construct full file path
        df.to_csv(filename, index=False)
        print(f"CSV file '{filename}' created successfully")

def write_to_csv_sgcar(data, filename, working_directory):  # Pass working directory as a parameter
    data_rows = []

    for (make_model, depreciation_value, price), price in data.items():
        data_rows.append({'Make-Model': make_model, 'Depreciation': depreciation_value, 'Price': price})
    
    df = pd.DataFrame(data_rows)
    df = df[df['Depreciation'] != "N.A"]
    df['Depreciation'] = pd.to_numeric(df['Depreciation'], errors='coerce')
    df['Price'] = pd.to_numeric(df['Price'], errors='coerce')
    df = df.dropna(subset=['Depreciation']).sort_values(by='Depreciation')

    if not df.empty:
        if not filename.endswith('.csv'):
            filename += '.csv'
        filename = os.path.join(working_directory, filename)  # Construct full file path
        df.to_csv(filename, index=False)
        print(f"CSV file '{filename}' created successfully")


def create_slide_with_table(prs, df, title):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    rows, cols = df.shape
    table = slide.shapes.add_table(rows=rows + 1, cols=cols, left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(1)).table

    for i, column in enumerate(df.columns):
        table.cell(0, i).text = column

    for i, row in enumerate(df.itertuples(), start=1):
        for j, value in enumerate(row[1:], start=0):
            table.cell(i, j).text = str(value)
    print(f"Slide with title '{title}' created successfully.")

def create_top_bottom_slides(prs, top3_dfs, bottom3_dfs):
    for top_df, bottom_df in zip(top3_dfs, bottom3_dfs):
        create_slide_with_table(prs, top_df[0], top_df[1])
        create_slide_with_table(prs, bottom_df[0], bottom_df[1])
    print("Slides with top and bottom data created successfully.")

print("Functions loaded successfully.")
